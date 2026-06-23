import io
import re
import tempfile
from pathlib import Path

from fastapi import FastAPI, File, HTTPException, UploadFile
from fastapi.responses import HTMLResponse, PlainTextResponse
from fastapi.staticfiles import StaticFiles

app = FastAPI(title="MDPreProcess")


# ── converters ────────────────────────────────────────────────────────────────

def pptx_to_md(data: bytes) -> str:
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN

    prs = Presentation(io.BytesIO(data))
    lines: list[str] = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_title = ""
        body_shapes = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape.shape_type == 13:  # picture
                continue
            try:
                ph = shape.placeholder_format
                if ph is not None and ph.idx == 0:
                    slide_title = shape.text_frame.text.strip()
                    continue
            except Exception:
                pass
            body_shapes.append(shape)

        heading = slide_title if slide_title else f"Slide {slide_num}"
        lines.append(f"## {heading}")

        for shape in body_shapes:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                level = para.level  # 0-based indent level
                prefix = "  " * level + "-"
                lines.append(f"{prefix} {text}")

        lines.append("")

    return "\n".join(lines).strip()


def pdf_to_md(data: bytes) -> str:
    import fitz  # pymupdf

    doc = fitz.open(stream=data, filetype="pdf")
    lines: list[str] = []

    for page_num, page in enumerate(doc, 1):
        blocks = page.get_text("dict")["blocks"]
        page_lines: list[tuple[float, str]] = []  # (font_size, text)

        for block in blocks:
            if block.get("type") != 0:  # 0 = text block
                continue
            for line in block.get("lines", []):
                spans = line.get("spans", [])
                if not spans:
                    continue
                text = "".join(s["text"] for s in spans).strip()
                if not text:
                    continue
                max_size = max(s["size"] for s in spans)
                page_lines.append((max_size, text))

        if not page_lines:
            continue

        # Determine heading threshold: top 10% font sizes on the page
        sizes = sorted({s for s, _ in page_lines}, reverse=True)
        heading_sizes = set(sizes[: max(1, len(sizes) // 5)])

        lines.append(f"## Page {page_num}")
        for size, text in page_lines:
            # Skip likely page numbers (short purely numeric strings)
            if re.fullmatch(r"\d{1,3}", text):
                continue
            if size in heading_sizes and len(page_lines) > 1:
                lines.append(f"### {text}")
            else:
                lines.append(f"- {text}")
        lines.append("")

    doc.close()
    return "\n".join(lines).strip()


# ── routes ────────────────────────────────────────────────────────────────────

@app.get("/", response_class=HTMLResponse)
async def index():
    html = Path("static/index.html").read_text()
    return HTMLResponse(content=html)


@app.post("/convert", response_class=PlainTextResponse)
async def convert(file: UploadFile = File(...)):
    ext = Path(file.filename).suffix.lower()
    if ext not in {".pptx", ".pdf"}:
        raise HTTPException(status_code=400, detail="Solo se aceptan archivos .pptx o .pdf")

    data = await file.read()
    try:
        if ext == ".pptx":
            md = pptx_to_md(data)
        else:
            md = pdf_to_md(data)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error al procesar el archivo: {e}")

    return PlainTextResponse(content=md, media_type="text/plain; charset=utf-8")


app.mount("/static", StaticFiles(directory="static"), name="static")
